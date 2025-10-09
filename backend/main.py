from flask import Flask, request, jsonify, render_template, session, redirect, url_for, send_from_directory
from flask_cors import CORS
from google import genai
from google.genai import types
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from flask_sqlalchemy import SQLAlchemy
from dotenv import load_dotenv
import os

load_dotenv()
API_KEY = os.getenv("GOOGLE_GEMINI_API_KEY")
SECRET_KEY = os.getenv("FLASK_SECRET_KEY")

# --- Flask Setup ---
app = Flask(__name__, template_folder="templates", static_folder="static")
CORS(app, supports_credentials=True, origins=["http://localhost:3000"])
app.secret_key = "SECRET_KEY"  # change for production

# --- SQLite Database Setup ---
app.config["SQLALCHEMY_DATABASE_URI"] = "sqlite:///quickiebot.db"
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False
db = SQLAlchemy(app)

# --- File Upload Config ---
UPLOAD_FOLDER = os.path.join("static", "uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx", "jpg", "jpeg", "png", "txt"}
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# --- Gemini API Setup ---

client = genai.Client(api_key=API_KEY)

# --- Database Models ---
class Admin(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password = db.Column(db.String(200), nullable=False)

class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password = db.Column(db.String(200), nullable=False)   

class Resource(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(200), nullable=False)
    description = db.Column(db.String(500))
    filepath = db.Column(db.String(500), nullable=False)
    thumbnail = db.Column(db.String(500))

# --- Helpers ---
def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

# Create DB tables
with app.app_context():
    db.create_all()

# --- Gemini Chat Endpoint ---
@app.route("/chat", methods=["POST"])
def chat():
    try:
        data = request.get_json()
        user_message = data.get("message", "")
        if not user_message.strip():
            return jsonify({"reply": "Please enter a valid message."}), 400
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=user_message,
            config=types.GenerateContentConfig(
                thinking_config=types.ThinkingConfig(thinking_budget=0)
            ),
        )
        return jsonify({"reply": response.text})
    except Exception as e:
        print("Error:", e)
        return jsonify({"reply": "Sorry, QuickieBot ran into an issue."}), 500
    

@app.route("/chat/file", methods=["POST"])
def chat_file():
    if "file" not in request.files:
        return jsonify({"reply": "No file uploaded."}), 400

    uploaded_file = request.files["file"]
    filename = uploaded_file.filename

    if filename == "":
        return jsonify({"reply": "No selected file."}), 400

    if not allowed_file(filename):
        return jsonify({"reply": "File type not allowed."}), 400

    content = ""

    try:
        ext = filename.rsplit(".", 1)[1].lower()

        # --- Handle text-based files ---
        if ext == "txt":
            content = uploaded_file.read().decode("utf-8", errors="ignore")

        elif ext == "pdf":
            import PyPDF2
            reader = PyPDF2.PdfReader(uploaded_file)
            for page in reader.pages:
                content += page.extract_text() or ""

        elif ext == "docx":
            import docx
            doc = docx.Document(uploaded_file)
            for para in doc.paragraphs:
                content += para.text + "\n"

        elif ext == "pptx":
            from pptx import Presentation
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"

        # --- Handle image files (jpg, jpeg, png) ---
        elif ext in {"jpg", "jpeg", "png"}:
            import base64
            import io
            from PIL import Image

            # Convert image to base64 for Gemini
            image = Image.open(uploaded_file)
            buffered = io.BytesIO()
            image.save(buffered, format=image.format)
            img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
            content = f"[Image uploaded: {filename}]"

            # Use Gemini’s multimodal capability
            response = client.models.generate_content(
                model="gemini-2.5-flash",
                contents=[
                    types.Part.from_text("Describe and summarize this image in educational context."),
                    types.Part.from_data(mime_type=f"image/{ext}", data=base64.b64decode(img_str)),
                ],
            )
            return jsonify({"reply": response.text})

        else:
            return jsonify({"reply": "Unsupported file type."}), 400

    except Exception as e:
        print("Error reading file:", e)
        return jsonify({"reply": f"Error reading file: {str(e)}"}), 500

    # --- Text-based summarization (default) ---
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=f"Summarize this content for students and provide examples:\n{content}",
            config=types.GenerateContentConfig(
                thinking_config=types.ThinkingConfig(thinking_budget=0)
            ),
        )
        return jsonify({"reply": response.text})
    except Exception as e:
        print("Gemini AI error:", e)
        return jsonify({"reply": "Error processing the file with AI."}), 500

# USER AUTH ROUTES

@app.route("/user/register", methods=["POST"])
def register_user():
    data = request.get_json()
    username = data.get("username")
    email = data.get("email")
    password = data.get("password")

    if not username or not email or not password:
        return jsonify({"error": "All fields are required."}), 400
    if User.query.filter_by(username=username).first():
        return jsonify({"error": "Username already exists."}), 400
    if User.query.filter_by(email=email).first():
        return jsonify({"error": "Email already registered."}), 400

    new_user = User(
        username=username,
        email=email,
        password=generate_password_hash(password),
    )
    db.session.add(new_user)
    db.session.commit()
    return jsonify({"message": "User registered successfully!"}), 201


@app.route("/user/login", methods=["POST"])
def login_user():
    data = request.get_json()
    username = data.get("username")
    password = data.get("password")

    user = User.query.filter_by(username=username).first()
    if not user or not check_password_hash(user.password, password):
        return jsonify({"error": "Invalid username or password."}), 401

    session["user_id"] = user.id
    session["username"] = user.username
    return jsonify({"message": f"Welcome, {user.username}!"})


@app.route("/user/logout", methods=["POST"])
def logout_user():
    session.clear()
    return jsonify({"message": "Logged out successfully."})

# --- Admin Register ---
@app.route("/admin/register", methods=["POST"])
def register_admin():
    data = request.get_json()
    username = data.get("username")
    email = data.get("email")
    password = data.get("password")
    if not username or not email or not password:
        return jsonify({"error": "All fields are required."}), 400
    if Admin.query.filter_by(username=username).first():
        return jsonify({"error": "Username already exists."}), 400
    new_admin = Admin(
        username=username,
        email=email,
        password=generate_password_hash(password),
    )
    db.session.add(new_admin)
    db.session.commit()
    return jsonify({"message": "Admin registered successfully!"}), 201


# --- Admin Login ---
@app.route("/admin/login", methods=["POST"])
def login_admin():
    data = request.get_json()
    username = data.get("username")
    password = data.get("password")
    admin = Admin.query.filter_by(username=username).first()
    if not admin or not check_password_hash(admin.password, password):
        return jsonify({"error": "Invalid username or password."}), 401
    session["admin_id"] = admin.id
    session["admin_username"] = admin.username
    return jsonify({"message": f"Welcome, {admin.username}!"})


# --- Admin Logout ---
@app.route("/admin/logout", methods=["POST"])
def logout_admin():
    session.clear()
    return jsonify({"message": "Logged out successfully."})


# --- File Upload CRUD ---
@app.route("/api/resources", methods=["GET"])
def list_resources():
    res_list = Resource.query.all()
    return jsonify([
        {"id": r.id, "filename": r.filename, "description": r.description, "filepath": r.filepath, "thumbnail": r.thumbnail}
        for r in res_list
    ])


@app.route("/api/resources", methods=["POST"])
def upload_resource():
    if "file" not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files["file"]
    thumbnail = request.files.get("thumbnail")
    description = request.form.get("description", "")
    if file.filename == "":
        return jsonify({"error": "No selected file"}), 400
    if not allowed_file(file.filename):
        return jsonify({"error": "File type not allowed"}), 400
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
    file.save(filepath)

    thumbnail_filename = None
    if thumbnail and thumbnail.filename != "":
        thumb_name = secure_filename(thumbnail.filename)
        thumb_path = os.path.join(app.config["UPLOAD_FOLDER"], thumb_name)
        thumbnail.save(thumb_path)
        thumbnail_filename = thumb_name

    resource = Resource(filename=filename, description=description, filepath=filepath, thumbnail=thumbnail_filename )
    db.session.add(resource)
    db.session.commit()
    return jsonify({
        "id": resource.id, 
        "filename": resource.filename, 
        "description": resource.description, 
        "filepath": resource.filepath,
        "thumbnail": resource.thumbnail
        }), 201

@app.route("/api/resources/<int:resource_id>", methods=["PUT"])
def update_resource(resource_id):
    resource = Resource.query.get(resource_id)
    if not resource:
        return jsonify({"error": "Resource not found"}), 404

    description = request.form.get("description")
    if description is not None:
        resource.description = description

    if "file" in request.files:
        file = request.files["file"]
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)

            # Remove old file
            try:
                os.remove(resource.filepath)
            except Exception:
                pass

            resource.filename = filename
            resource.filepath = filepath

    db.session.commit()
    return jsonify({
        "id": resource.id,
        "filename": resource.filename,
        "description": resource.description,
        "filepath": resource.filepath
    })

@app.route("/api/resources/<int:resource_id>", methods=["DELETE"])
def delete_resource(resource_id):
    resource = Resource.query.get(resource_id)
    if not resource:
        return jsonify({"error": "Resource not found"}), 404
    try:
        os.remove(resource.filepath)
    except Exception:
        pass
    db.session.delete(resource)
    db.session.commit()
    return jsonify({"message": "Resource deleted"})


@app.route("/uploads/<filename>")
def uploaded_file(filename):
    return send_from_directory(app.config["UPLOAD_FOLDER"], filename)


# --- Admin Panel ---
@app.route("/admin")
def admin_panel():
    if "admin_id" not in session:
        return redirect(url_for("login_page"))
    return render_template(
        "admin.html",
        username=session["admin_username"],
        resources=Resource.query.all(),
    )


# --- Login + Register Pages ---
@app.route("/admin/login", methods=["GET"])
def login_page():
    return render_template("login.html")


@app.route("/admin/register", methods=["GET"])
def register_page():
    return render_template("register.html")


@app.route("/")
def home():
    return "Backend running"

with app.app_context():
    db.create_all()  # Creates tables if they don't exist

if __name__ == "__main__":
    debug_mode = os.getenv("FLASK_DEBUG", "false").lower() == "true"
    port = int(os.getenv("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=debug_mode)